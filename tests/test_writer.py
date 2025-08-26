import pytest
import sys
import os
from io import BytesIO
import tempfile

# Add the backend directory to the Python path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'backend'))

# Note: These tests require python-pptx to be installed
# They will be skipped if the library is not available

def test_build_presentation_basic():
    """Test basic presentation building."""
    try:
        from backend.pptx_engine.slide_writer import build_presentation
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx not installed")
    
    # Create a minimal template
    template_prs = Presentation()
    template_path = None
    
    try:
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            template_prs.save(tmp)
            template_path = tmp.name
        
        # Test plan
        plan = {
            "slides": [
                {
                    "title": "Test Slide",
                    "layout_hint": "bullets",
                    "bullets": ["Point 1", "Point 2"],
                    "notes": "Test notes"
                }
            ]
        }
        
        # Mock style context
        style_ctx = {
            "theme": {"font_family": "Arial"},
            "images": [],
            "layouts": {}
        }
        
        # Build presentation
        result_bytes = build_presentation(plan, style_ctx, template_path)
        
        # Verify result
        assert isinstance(result_bytes, bytes)
        assert len(result_bytes) > 0
        
        # Verify it's a valid PPTX by opening it
        result_prs = Presentation(BytesIO(result_bytes))
        assert len(result_prs.slides) == 1
        
    finally:
        if template_path and os.path.exists(template_path):
            os.remove(template_path)

def test_build_presentation_multiple_slides():
    """Test building presentation with multiple slides."""
    try:
        from backend.pptx_engine.slide_writer import build_presentation
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx not installed")
    
    # Create a minimal template
    template_prs = Presentation()
    template_path = None
    
    try:
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            template_prs.save(tmp)
            template_path = tmp.name
        
        # Test plan with multiple slides
        plan = {
            "slides": [
                {
                    "title": "Introduction",
                    "layout_hint": "bullets",
                    "bullets": ["Welcome", "Overview", "Agenda"]
                },
                {
                    "title": "Main Content",
                    "layout_hint": "section",
                    "bullets": ["Key Point 1", "Key Point 2"]
                },
                {
                    "title": "Conclusion",
                    "layout_hint": "bullets", 
                    "bullets": ["Summary", "Next Steps"],
                    "notes": "Thank the audience"
                }
            ]
        }
        
        style_ctx = {
            "theme": {"font_family": "Calibri"},
            "images": [],
            "layouts": {}
        }
        
        result_bytes = build_presentation(plan, style_ctx, template_path)
        
        # Verify result
        assert isinstance(result_bytes, bytes)
        assert len(result_bytes) > 0
        
        # Verify slide count
        result_prs = Presentation(BytesIO(result_bytes))
        assert len(result_prs.slides) == 3
        
    finally:
        if template_path and os.path.exists(template_path):
            os.remove(template_path)

def test_build_presentation_with_long_content():
    """Test that long bullet points are handled properly."""
    try:
        from backend.pptx_engine.slide_writer import build_presentation
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx not installed")
    
    template_prs = Presentation()
    template_path = None
    
    try:
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            template_prs.save(tmp)
            template_path = tmp.name
        
        # Test with many long bullet points
        long_bullets = [f"This is a very long bullet point number {i} that should be handled gracefully by the system" for i in range(10)]
        
        plan = {
            "slides": [
                {
                    "title": "Test with Long Content",
                    "layout_hint": "bullets",
                    "bullets": long_bullets
                }
            ]
        }
        
        style_ctx = {
            "theme": {"font_family": "Arial"},
            "images": [],
            "layouts": {}
        }
        
        result_bytes = build_presentation(plan, style_ctx, template_path)
        
        # Should not crash and should produce valid output
        assert isinstance(result_bytes, bytes)
        assert len(result_bytes) > 0
        
        result_prs = Presentation(BytesIO(result_bytes))
        assert len(result_prs.slides) == 1
        
    finally:
        if template_path and os.path.exists(template_path):
            os.remove(template_path)

def test_layout_selection():
    """Test the layout selection logic."""
    try:
        from backend.pptx_engine.slide_writer import _pick_layout
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx not installed")
    
    prs = Presentation()
    
    # Test different layout hints
    layout1 = _pick_layout(prs, "bullets")
    layout2 = _pick_layout(prs, "section")
    layout3 = _pick_layout(prs, "two-column")
    
    # Should return layout objects (exact layout depends on template)
    assert layout1 is not None
    assert layout2 is not None  
    assert layout3 is not None

if __name__ == "__main__":
    pytest.main([__file__])
