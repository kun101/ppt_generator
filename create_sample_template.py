#!/usr/bin/env python3
"""
Create a sample PowerPoint template for testing purposes.
"""

try:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    
    # Create a new presentation
    prs = Presentation()
    
    # Add a title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Sample Template"
    subtitle.text = "Created for Text-to-PowerPoint Generator Testing"
    
    # Add a content slide
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Sample Content Slide"
    content.text = "This is a sample bullet point\nThis is another bullet point"
    
    # Save the template
    prs.save('sample_template.pptx')
    print("✅ Sample template created: sample_template.pptx")
    
except ImportError:
    print("❌ python-pptx not available. Cannot create sample template.")
except Exception as e:
    print(f"❌ Error creating template: {e}")
