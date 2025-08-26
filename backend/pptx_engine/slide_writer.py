from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
import tempfile
import os
import re
from io import BytesIO

def build_presentation(plan, style_ctx, template_path: str) -> bytes:
    """Build a PowerPoint presentation from the slide plan and template."""
    # Create new presentation from template
    prs = Presentation(template_path)

    # Ensure style context has slide dimension metadata (backward compatibility)
    if "theme" in style_ctx and "slide_width_emu" not in style_ctx["theme"]:
        try:
            style_ctx["theme"].update({
                "slide_width_emu": prs.slide_width,
                "slide_height_emu": prs.slide_height,
                "slide_width_inches": float(prs.slide_width) / 914400.0,
                "slide_height_inches": float(prs.slide_height) / 914400.0,
                "aspect_ratio": round((float(prs.slide_width) / float(prs.slide_height)), 4) if prs.slide_height else 1.0
            })
        except Exception:
            pass
    
    # Clear existing slides but keep the template structure
    slide_idxs = list(range(len(prs.slides) - 1, -1, -1))
    for i in slide_idxs:
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    
    # Prepare image pool for intelligent reuse
    image_pool = style_ctx["images"]
    image_usage_count = 0
    
    # Generate slides from plan
    for i, slide_def in enumerate(plan["slides"]):
        # Choose layout based on hint
        layout = _pick_layout(prs, slide_def.get("layout_hint", "bullets"))
        slide = prs.slides.add_slide(layout)
        
        # Get layout info for proper font styling
        layout_idx = _get_layout_index(prs, layout)
        layout_info = style_ctx["theme"]["layout_info"].get(layout_idx, {})
        
        # Fill placeholders (text first, then images) and learn if an image placeholder was used
        fill_result = _fill_all_placeholders(slide, slide_def, style_ctx, layout_info, image_pool)
        used_image_placeholder = fill_result.get("used_image_placeholder", False) if isinstance(fill_result, dict) else False

        # COMPLETELY SKIP IMAGE HANDLING FOR NOW - TEXT PRIORITY ONLY
        # TODO: Add intelligent image placement after all text is properly positioned
        
        # Add speaker notes
        notes = slide_def.get("notes", "")
        if notes and slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = notes
    
    # Save to bytes
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()

def _get_layout_index(prs, target_layout):
    """Get the index of a layout in the presentation."""
    for i, layout in enumerate(prs.slide_layouts):
        if layout == target_layout:
            return i
    return 1  # Default to content layout

def _should_add_image(slide_def, slide_index):
    """Determine if this slide should get an image - much more conservative."""
    layout_hint = slide_def.get("layout_hint", "").lower()
    
    # Only add images to explicitly image-focused layouts
    if any(keyword in layout_hint for keyword in ["image", "picture", "visual", "photo"]):
        return True
    
    # Very rarely add images to other slides to avoid overlap
    return False

def _populate_title(title_shape, title_text, layout_info):
    """Populate title with proper template formatting."""
    if not hasattr(title_shape, 'text_frame'):
        return
        
    text_frame = title_shape.text_frame
    text_frame.clear()
    
    # Get title placeholder font info
    title_font_info = None
    for placeholder in layout_info.get("placeholders", []):
        if placeholder.get("type") == 1:  # Title placeholder
            title_font_info = placeholder.get("font")
            break
    
    # Apply formatted text with template font
    paragraph = text_frame.paragraphs[0]
    _apply_formatted_text(paragraph, title_text, title_font_info, is_title=True)

def _populate_bullets(shape, bullets, layout_info):
    """Populate a text shape with bullet points using template formatting."""
    if not hasattr(shape, 'text_frame'):
        return
    
    # Get body placeholder font info
    body_font_info = None
    for placeholder in layout_info.get("placeholders", []):
        if placeholder.get("type") == 2:  # Body placeholder
            body_font_info = placeholder.get("font")
            break
    
    text_frame = shape.text_frame
    text_frame.clear()
    
    for i, bullet in enumerate(bullets[:6]):  # Limit to 6 bullets
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        _apply_formatted_text(p, str(bullet), body_font_info, is_title=False)
        p.level = 0

def _pick_layout(prs, hint: str):
    """Select an appropriate layout based on the hint, prioritizing layouts with image placeholders."""
    hint = hint.lower()
    layouts = prs.slide_layouts
    
    # First, try to find layouts with image placeholders if hint suggests images
    if any(word in hint for word in ["image", "picture", "photo", "visual"]):
        for i, layout in enumerate(layouts):
            if _layout_has_image_placeholder(layout):
                return layout
    
    # Map hints to layout indices (varies by template)
    if "section" in hint and len(layouts) > 2:
        return layouts[2]  # Section header layout
    elif "two" in hint and len(layouts) > 3:
        return layouts[3]  # Two-column layout
    elif "title" in hint and len(layouts) > 0:
        return layouts[0]  # Title slide
    else:
        # Default to title+content layout, but prefer one with image placeholder
        for i in range(1, min(len(layouts), 4)):
            if _layout_has_image_placeholder(layouts[i]):
                return layouts[i]
        
        # Fallback to standard content layout
        return layouts[1] if len(layouts) > 1 else layouts[0]

def _layout_has_image_placeholder(layout):
    """Check if a layout has image placeholders."""
    for shape in layout.shapes:
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            # Check for image placeholder types
            if placeholder_type in [18, 19, 20, 7]:  # Picture, clip art, media, object
                return True
    return False

def _fill_all_placeholders(slide, slide_def, style_ctx, layout_info, image_pool):
    """Fill placeholders systematically - text content first, then images in designated placeholders only.

    Returns: {"used_image_placeholder": bool}
    """
    title_text = slide_def.get("title", "")
    bullets = slide_def.get("bullets", [])
    layout_hint = slide_def.get("layout_hint", "").lower()

    filled_placeholders = set()
    used_image_placeholder = False
    content_usage_tracker = {
        "title_used": False,
        "bullets_used": set(),  # Track which bullets have been used
        "main_content_filled": False
    }

    print(f"DEBUG: Starting placeholder fill for slide with title: '{title_text}' and {len(bullets)} bullets")
    
    # First pass: Fill text placeholders with strict content assignment
    for shape in slide.shapes:
        if not (hasattr(shape, 'is_placeholder') and shape.is_placeholder):
            continue
        
        placeholder_type = shape.placeholder_format.type
        placeholder_idx = shape.placeholder_format.idx
        
        print(f"DEBUG: Found placeholder - Type: {placeholder_type}, Index: {placeholder_idx}, Has text_frame: {hasattr(shape, 'text_frame')}")

        # Skip image placeholders in first pass
        if placeholder_type in [18, 19, 20, 7]:  # Picture, clip art, media, object
            continue

        # Title placeholder (Type 1) - use title text only once
        if placeholder_type == 1 and placeholder_idx not in filled_placeholders:
            if title_text and not content_usage_tracker["title_used"]:
                print(f"DEBUG: Filling title placeholder {placeholder_idx} with title text")
                _fill_text_placeholder_safely(shape, title_text, layout_info, is_title=True)
                filled_placeholders.add(placeholder_idx)
                content_usage_tracker["title_used"] = True
            continue

        # Main body placeholder (Type 2) - use all bullets once
        if placeholder_type == 2 and placeholder_idx not in filled_placeholders:
            if bullets and not content_usage_tracker["main_content_filled"]:
                print(f"DEBUG: Filling main body placeholder {placeholder_idx} with {len(bullets)} bullets")
                _fill_body_placeholder_with_expansion(shape, bullets, layout_info)
                filled_placeholders.add(placeholder_idx)
                content_usage_tracker["main_content_filled"] = True
                # Mark all bullets as used in main content
                for i in range(len(bullets)):
                    content_usage_tracker["bullets_used"].add(i)
            continue

        # Subtitle placeholder (Type 3) - use specific subtitle content
        if placeholder_type == 3 and placeholder_idx not in filled_placeholders:
            subtitle_text = _generate_subtitle_content(slide_def, layout_hint, content_usage_tracker)
            if subtitle_text:
                print(f"DEBUG: Filling subtitle placeholder {placeholder_idx} with subtitle")
                _fill_text_placeholder_safely(shape, subtitle_text, layout_info, is_title=False)
                filled_placeholders.add(placeholder_idx)
            continue

        # Content placeholder (Type 8) - use remaining content
        if placeholder_type == 8 and hasattr(shape, 'text_frame') and shape.text_frame is not None and placeholder_idx not in filled_placeholders:
            remaining_content = _get_remaining_content(slide_def, content_usage_tracker)
            if remaining_content:
                limited_content = _limit_content_for_placeholder(remaining_content, shape)
                print(f"DEBUG: Filling content placeholder {placeholder_idx} (type 8) with remaining content")
                _fill_text_placeholder_safely(shape, limited_content, layout_info, is_title=False)
                filled_placeholders.add(placeholder_idx)
            continue

        # Other text placeholders - use very specific content to avoid duplication
        if hasattr(shape, 'text_frame') and shape.text_frame is not None and placeholder_idx not in filled_placeholders:
            specific_content = _generate_specific_placeholder_content(slide_def, placeholder_type, layout_hint, content_usage_tracker)
            if specific_content:
                limited_content = _limit_content_for_placeholder(specific_content, shape)
                print(f"DEBUG: Filling text placeholder {placeholder_idx} (type {placeholder_type}) with specific content")
                _fill_text_placeholder_safely(shape, limited_content, layout_info, is_title=False)
                filled_placeholders.add(placeholder_idx)

    # Second pass: Fill image placeholders ONLY if images are available and layout supports them
    if image_pool and _should_use_images_for_layout(layout_hint):
        for shape in slide.shapes:
            if not (hasattr(shape, 'is_placeholder') and shape.is_placeholder):
                continue
            
            placeholder_type = shape.placeholder_format.type
            placeholder_idx = shape.placeholder_format.idx
            
            # Only fill image placeholders, maintain their exact size and position
            if placeholder_type in [18, 19, 20, 7] and placeholder_idx not in filled_placeholders:  # Image placeholder types
                print(f"DEBUG: Found image placeholder {placeholder_idx} (type {placeholder_type})")
                if _fill_image_placeholder_safely(slide, shape, image_pool, len(filled_placeholders)):
                    filled_placeholders.add(placeholder_idx)
                    used_image_placeholder = True
                    print(f"DEBUG: Successfully filled image placeholder {placeholder_idx}")

    print(f"DEBUG: Placeholder filling complete. Filled: {len(filled_placeholders)} placeholders")
    return {"used_image_placeholder": used_image_placeholder}

def _should_use_image_placeholder(layout_hint):
    """Determine if we should use image placeholders based on layout hint."""
    # Only use image placeholders for explicitly image-focused layouts
    return any(keyword in layout_hint for keyword in ["image", "picture", "visual", "photo"])

def _get_slide_area():
    """Get approximate slide area for sizing calculations."""
    try:
        return Inches(10) * Inches(7.5)  # Standard 16:9 slide
    except Exception:
        return 914400 * 685800  # EMU units fallback

def _fill_text_placeholder_safely(shape, text, layout_info, is_title=False):
    """Fill text placeholder with ultra-conservative sizing to prevent any overflow."""
    if not hasattr(shape, 'text_frame') or shape.text_frame is None:
        return
    
    # Get font info from layout
    font_info = None
    placeholder_type = 1 if is_title else 2
    for placeholder in layout_info.get("placeholders", []):
        if placeholder.get("type") == placeholder_type:
            font_info = placeholder.get("font")
            break
    
    text_frame = shape.text_frame
    text_frame.clear()
    
    # Configure text frame for better text flow
    text_frame.word_wrap = True
    text_frame.auto_size = None  # Disable auto-resize to prevent expansion
    
    if text.strip():
        paragraph = text_frame.paragraphs[0]
        
        # Ultra-conservative font sizing based on shape size and text length
        safe_font_info = _calculate_ultra_safe_font(font_info, text, shape, is_title)
        _apply_formatted_text(paragraph, text, safe_font_info, is_title)
        
        print(f"DEBUG: Filled {'title' if is_title else 'text'} placeholder with font size {safe_font_info.get('size', 'default')}")

def _calculate_ultra_safe_font(font_info, text, shape, is_title):
    """Calculate ultra-conservative font size to absolutely prevent overflow."""
    if not font_info:
        base_size = 24 if is_title else 18
        font_info = {"name": "Calibri", "size": base_size, "bold": False, "italic": False}
    else:
        font_info = font_info.copy()  # Don't modify original
    
    # Get shape dimensions in a safer way
    try:
        shape_width_inches = shape.width / 914400
        shape_height_inches = shape.height / 914400
        shape_area = shape_width_inches * shape_height_inches
    except:
        shape_area = 10  # Fallback assumption
    
    text_length = len(text)
    base_size = font_info.get("size", 24 if is_title else 18)
    
    # Ultra-aggressive size reduction based on multiple factors
    size_reduction = 1.0
    
    # Factor 1: Text length
    if text_length > 150:
        size_reduction *= 0.5  # Very long text
    elif text_length > 100:
        size_reduction *= 0.6
    elif text_length > 50:
        size_reduction *= 0.7
    elif text_length > 30:
        size_reduction *= 0.8
    
    # Factor 2: Shape size
    if shape_area < 5:  # Very small shape
        size_reduction *= 0.5
    elif shape_area < 10:
        size_reduction *= 0.6
    elif shape_area < 20:
        size_reduction *= 0.8
    
    # Factor 3: Title vs body text
    if is_title and text_length > 50:
        size_reduction *= 0.7  # Long titles need to be smaller
    
    # Apply the reduction
    new_size = max(int(base_size * size_reduction), 8 if is_title else 6)  # Minimum readable size
    font_info["size"] = new_size
    
    print(f"DEBUG: Ultra-safe font calculation: text_len={text_length}, shape_area={shape_area:.2f}, reduction={size_reduction:.2f}, final_size={new_size}")
    
    return font_info

def _fill_body_placeholder_with_expansion(shape, bullets, layout_info):
    """Fill body placeholder with conservative sizing to prevent overflow."""
    if not hasattr(shape, 'text_frame') or shape.text_frame is None:
        return
    
    # Get body placeholder font info
    body_font_info = None
    for placeholder in layout_info.get("placeholders", []):
        if placeholder.get("type") == 2:  # Body placeholder
            body_font_info = placeholder.get("font")
            break
    
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    # CRITICAL: Don't expand placeholders - work within existing bounds
    # Use the original shape size to prevent overflow
    available_height = shape.height - Inches(0.4)  # Conservative margin
    
    # Calculate how many bullets can safely fit
    max_bullets = _calculate_safe_bullets_for_shape(shape, body_font_info, bullets)
    safe_bullets = bullets[:max_bullets]
    
    print(f"DEBUG: Shape height: {shape.height}, Available height: {available_height}, Max bullets: {max_bullets}")
    
    for i, bullet_text in enumerate(safe_bullets):
        if i == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        
        paragraph.level = 0
        # Use appropriately sized font for the available space
        safe_font_info = _calculate_safe_font_size(body_font_info, len(safe_bullets), shape)
        _apply_formatted_text(paragraph, str(bullet_text), safe_font_info, is_title=False)

def _calculate_safe_bullets_for_shape(shape, font_info, bullets):
    """Calculate maximum bullets that will definitely fit without overflow."""
    # Start with very conservative assumptions
    font_size = font_info.get("size", 18) if font_info else 18
    
    # Conservative line height calculation
    line_height_emu = int(font_size * 1.5 * 12700)  # Convert to EMU with extra spacing
    
    # Available height with generous margins
    margin_emu = int(Inches(0.4))  # Top and bottom margins
    available_height_emu = shape.height - margin_emu
    
    # Estimate text height for each bullet
    total_height_needed = 0
    max_bullets = 0
    
    for i, bullet in enumerate(bullets):
        # Estimate lines needed for this bullet (conservative: assume 40 chars per line)
        estimated_lines = max(1, len(str(bullet)) // 40 + 1)
        bullet_height = estimated_lines * line_height_emu
        
        if total_height_needed + bullet_height <= available_height_emu:
            total_height_needed += bullet_height
            max_bullets += 1
        else:
            break
    
    # Be extra conservative - reduce by 1 if we have more than 1 bullet
    if max_bullets > 1:
        max_bullets = max(1, max_bullets - 1)
    
    print(f"DEBUG: Font size: {font_size}, Line height: {line_height_emu}, Available: {available_height_emu}, Safe bullets: {max_bullets}")
    
    return max_bullets

def _calculate_safe_font_size(font_info, bullet_count, shape):
    """Calculate a safe font size that will fit in the available space."""
    if not font_info:
        base_size = 18
        font_info = {"name": "Calibri", "size": base_size, "bold": False, "italic": False}
    else:
        font_info = font_info.copy()
    
    base_size = font_info.get("size", 18)
    
    # Reduce font size based on number of bullets and shape size
    shape_area_inches = (shape.width / 914400) * (shape.height / 914400)
    
    # Very aggressive font size reduction to prevent overflow
    if bullet_count > 5 or shape_area_inches < 20:  # Small shape
        font_info["size"] = max(int(base_size * 0.6), 10)
    elif bullet_count > 3 or shape_area_inches < 30:
        font_info["size"] = max(int(base_size * 0.7), 12)
    elif bullet_count > 2:
        font_info["size"] = max(int(base_size * 0.8), 14)
    else:
        font_info["size"] = max(int(base_size * 0.9), 16)
    
    print(f"DEBUG: Adjusted font size from {base_size} to {font_info['size']} for {bullet_count} bullets")
    
    return font_info

def _adjust_font_for_content(font_info, text, shape, is_title):
    """Adjust font size based on content length and shape size."""
    if not font_info:
        base_size = 24 if is_title else 18
        font_info = {"name": "Calibri", "size": base_size, "bold": False, "italic": False}
    else:
        font_info = font_info.copy()  # Don't modify original
    
    # Calculate scaling factor based on text length and shape size
    text_length = len(text)
    shape_area = shape.width * shape.height
    
    # Reduce font size for long text in small shapes
    if text_length > 100 and shape_area < Inches(4) * Inches(2):
        font_info["size"] = max(int(font_info.get("size", 18) * 0.8), 12)
    elif text_length > 200:
        font_info["size"] = max(int(font_info.get("size", 18) * 0.7), 10)
    
    return font_info

def _estimate_text_height(bullets, font_info):
    """Estimate height needed for bullet text."""
    if not bullets:
        return Inches(1)
    
    font_size = font_info.get("size", 18) if font_info else 18
    line_height = Pt(font_size * 1.2)  # Typical line spacing
    
    # Estimate lines per bullet (assuming ~50 chars per line at normal sizes)
    total_lines = 0
    for bullet in bullets:
        bullet_lines = max(1, len(str(bullet)) // 50)
        total_lines += bullet_lines
    
    return line_height * total_lines + Inches(0.5)  # Add padding

def _calculate_max_bullets_for_shape(shape, font_info):
    """Calculate maximum bullets that fit in shape without overflow."""
    font_size = font_info.get("size", 18) if font_info else 18
    line_height = Pt(font_size * 1.2)
    
    # Calculate available lines in shape
    available_height = shape.height - Inches(0.2)  # Margin
    max_lines = int(available_height / line_height)
    
    # Assume average 1.5 lines per bullet
    return max(1, int(max_lines / 1.5))

def _adjust_font_for_bullets(font_info, bullet_count):
    """Adjust font size based on number of bullets."""
    if not font_info:
        font_info = {"name": "Calibri", "size": 18, "bold": False, "italic": False}
    else:
        font_info = font_info.copy()
    
    # Reduce font size for many bullets
    if bullet_count > 6:
        font_info["size"] = max(int(font_info.get("size", 18) * 0.8), 12)
    elif bullet_count > 4:
        font_info["size"] = max(int(font_info.get("size", 18) * 0.9), 14)
    
    return font_info

def _limit_content_for_placeholder(content, shape):
    """Limit content length based on placeholder size to prevent overflow."""
    try:
        shape_area_inches = (shape.width / 914400) * (shape.height / 914400)
    except:
        shape_area_inches = 10
    
    # Calculate safe character limit based on shape size
    if shape_area_inches < 5:
        max_chars = 50  # Very small placeholder
    elif shape_area_inches < 10:
        max_chars = 100  # Small placeholder
    elif shape_area_inches < 20:
        max_chars = 200  # Medium placeholder
    else:
        max_chars = 400  # Large placeholder
    
    if len(content) > max_chars:
        # Truncate content and add ellipsis
        truncated = content[:max_chars-3] + "..."
        print(f"DEBUG: Truncated content from {len(content)} to {len(truncated)} chars for shape area {shape_area_inches:.2f}")
        return truncated
    
    return content

def _get_remaining_content(slide_def, content_usage_tracker):
    """Get content that hasn't been used yet."""
    bullets = slide_def.get("bullets", [])
    
    # If main content hasn't been filled, return bullets
    if not content_usage_tracker["main_content_filled"]:
        return "\n".join([f"• {bullet}" for bullet in bullets])
    
    # Find unused bullets
    unused_bullets = []
    for i, bullet in enumerate(bullets):
        if i not in content_usage_tracker["bullets_used"]:
            unused_bullets.append(bullet)
    
    if unused_bullets:
        return "\n".join([f"• {bullet}" for bullet in unused_bullets[:3]])  # Max 3 for secondary placeholders
    
    return ""

def _generate_specific_placeholder_content(slide_def, placeholder_type, layout_hint, content_usage_tracker):
    """Generate very specific content for placeholders to avoid duplication."""
    
    # For footer-type placeholders (Type 13, 14, 15, 16)
    if placeholder_type in [13, 14, 15, 16]:
        if "financial" in layout_hint:
            return "Key financial insights"
        elif "team" in layout_hint:
            return "Leadership team"
        elif "product" in layout_hint:
            return "Product overview"
        else:
            return "Supporting information"
    
    # For any other placeholder types, only use if we have unused content
    bullets = slide_def.get("bullets", [])
    unused_bullets = []
    for i, bullet in enumerate(bullets):
        if i not in content_usage_tracker["bullets_used"]:
            unused_bullets.append(bullet)
    
    # Only use 1 bullet for miscellaneous placeholders
    if unused_bullets and placeholder_type not in [1, 2, 3, 8]:
        content_usage_tracker["bullets_used"].add(bullets.index(unused_bullets[0]))
        return unused_bullets[0]
    
    return ""

def _should_use_images_for_layout(layout_hint):
    """Determine if this layout should use images based on layout hint."""
    # Only use images for layouts that explicitly mention images or visual content
    image_keywords = ["image", "picture", "visual", "photo", "media", "graphic"]
    return any(keyword in layout_hint.lower() for keyword in image_keywords)

def _fill_image_placeholder_safely(slide, shape, image_pool, usage_count):
    """Fill image placeholder using exact template positioning and sizing."""
    if not image_pool:
        return False
    
    try:
        # Select image based on usage pattern
        image_idx = usage_count % len(image_pool)
        image_data = image_pool[image_idx]
        
        # Create temporary file for the image
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
            tmp_file.write(image_data["blob"])
            tmp_path = tmp_file.name
        
        # CRITICAL: Store placeholder's exact position and size before removal
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        
        # Remove the placeholder
        placeholder_element = shape.element
        placeholder_element.getparent().remove(placeholder_element)
        
        # Add image with EXACT same position and size as placeholder using slide's shapes collection
        slide.shapes.add_picture(tmp_path, left, top, width, height)
        
        print(f"DEBUG: Image placed in placeholder at position ({left}, {top}) with size ({width}, {height})")
        return True
        
    except Exception as e:
        print(f"DEBUG: Failed to fill image placeholder: {e}")
        return False
    finally:
        try:
            os.remove(tmp_path)
        except:
            pass

def _generate_subtitle_content(slide_def, layout_hint, content_usage_tracker):
    """Generate appropriate subtitle content that hasn't been used."""
    bullets = slide_def.get("bullets", [])
    
    # For title slides, use first bullet as subtitle if title hasn't been used for main content
    if "title" in layout_hint and bullets and not content_usage_tracker["main_content_filled"]:
        # Use first bullet for subtitle and mark it as used
        content_usage_tracker["bullets_used"].add(0)
        return bullets[0]
    
    # For section slides, create a summary
    if "section" in layout_hint:
        return "Key insights and developments"
    
    # Generic subtitle based on context
    if bullets and not content_usage_tracker["main_content_filled"]:
        return "Overview and key points"
    
    return ""

def _apply_text_with_template_formatting(shape, text, layout_info, is_title=False):
    """Apply text with enhanced template-based formatting."""
    if not hasattr(shape, 'text_frame') or shape.text_frame is None:
        return
    
    # Get font info from layout
    font_info = None
    placeholder_type = 1 if is_title else 2
    for placeholder in layout_info.get("placeholders", []):
        if placeholder.get("type") == placeholder_type:
            font_info = placeholder.get("font")
            break
    
    # Clear existing content
    text_frame = shape.text_frame
    text_frame.clear()
    
    # Add the text with formatting
    if text.strip():
        paragraph = text_frame.paragraphs[0]
        _apply_formatted_text(paragraph, text, font_info, is_title)

def _populate_bullets_enhanced(shape, bullets, layout_info):
    """Enhanced bullet population with template formatting."""
    if not hasattr(shape, 'text_frame') or shape.text_frame is None:
        return
    
    # Get body placeholder font info
    body_font_info = None
    for placeholder in layout_info.get("placeholders", []):
        if placeholder.get("type") == 2:  # Body placeholder
            body_font_info = placeholder.get("font")
            break
    
    text_frame = shape.text_frame
    text_frame.clear()
    
    for i, bullet_text in enumerate(bullets[:6]):  # Limit to 6 bullets
        if i == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        
        paragraph.level = 0
        _apply_formatted_text(paragraph, str(bullet_text), body_font_info, is_title=False)

def _find_body_placeholder(slide):
    """Find the main content placeholder on a slide."""
    for shape in slide.shapes:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            # Type 2 is typically the body/content placeholder
            if ph_type == 2:  # PP_PLACEHOLDER.BODY
                return shape
    
    # Fallback: look for any text placeholder that's not the title
    for shape in slide.shapes:
        if (shape.is_placeholder and 
            hasattr(shape, 'text_frame') and 
            shape != getattr(slide.shapes, 'title', None)):
            return shape
    
    return None

def _apply_formatted_text(paragraph, text, font_info, is_title=False):
    """Apply Markdown formatting to a paragraph with template font info."""
    # Clear existing text
    paragraph.text = ""
    
    # Parse markdown formatting in the text
    parts = _parse_markdown_formatting(text)
    
    # Get font info from template or use defaults
    if font_info:
        font_name = font_info.get("name", "Calibri")
        font_size = Pt(font_info.get("size", 24 if is_title else 18))
        default_bold = font_info.get("bold", False)
        default_italic = font_info.get("italic", False)
    else:
        font_name = "Calibri"
        font_size = Pt(24 if is_title else 18)
        default_bold = False
        default_italic = False
    
    for part_text, is_bold, is_italic in parts:
        if not part_text:
            continue
            
        # Add a text run to the paragraph
        run = paragraph.add_run()
        run.text = part_text
        
        # Apply base font styling from template
        run.font.name = font_name
        run.font.size = font_size
        
        # Apply formatting (template defaults + markdown overrides)
        run.font.bold = default_bold or is_bold
        run.font.italic = default_italic or is_italic

def _add_intelligent_image(slide, image_pool, usage_count, layout_hint):
    """Add small, non-intrusive image that doesn't interfere with text."""
    if not image_pool:
        return
    
    # Select image based on usage pattern
    image_idx = usage_count % len(image_pool)
    image_data = image_pool[image_idx]
    
    try:
        # Create temporary file for the image
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
            tmp_file.write(image_data["blob"])
            tmp_path = tmp_file.name
        
        # First check for unused image placeholders
        image_placeholder = _find_image_placeholder(slide)
        
        if image_placeholder:
            # Use placeholder but make it smaller if too big
            _replace_placeholder_with_image_safely(slide, image_placeholder, tmp_path)
        else:
            # Only add non-placeholder image if layout explicitly requests it
            if any(keyword in layout_hint.lower() for keyword in ["image", "visual"]):
                placement = _calculate_conservative_image_placement(slide, image_data)
                if placement:
                    slide.shapes.add_picture(
                        tmp_path, 
                        placement["left"], 
                        placement["top"], 
                        placement["width"], 
                        placement["height"]
                    )
        
    except Exception as e:
        print(f"Failed to add image: {e}")
    finally:
        try:
            os.remove(tmp_path)
        except:
            pass

def _replace_placeholder_with_image_safely(slide, placeholder, image_path):
    """Replace placeholder with image, but limit size to prevent text overlap."""
    # Get placeholder properties
    left = placeholder.left
    top = placeholder.top
    original_width = placeholder.width
    original_height = placeholder.height
    
    # Limit image size to reasonable bounds - ensure integers
    max_width = int(min(original_width, Inches(3)))  # Max 3 inches wide
    max_height = int(min(original_height, Inches(2.5)))  # Max 2.5 inches tall
    
    # Remove the placeholder
    placeholder_element = placeholder.element
    placeholder_element.getparent().remove(placeholder_element)
    
    # Add image with limited size
    slide.shapes.add_picture(image_path, left, top, max_width, max_height)

def _calculate_conservative_image_placement(slide, image_data):
    """Calculate very conservative image placement that avoids all text."""
    # Get slide dimensions
    try:
        slide_width = slide.part.slide_width
        slide_height = slide.part.slide_height
    except Exception:
        slide_width = Inches(10)
        slide_height = Inches(7.5)
    
    # Find all text areas to avoid
    text_areas = []
    for shape in slide.shapes:
        if hasattr(shape, 'left') and hasattr(shape, 'top'):
            if hasattr(shape, 'text_frame') or shape.is_placeholder:
                # Add buffer around text areas
                buffer = Inches(0.3)
                text_areas.append({
                    "left": shape.left - buffer,
                    "top": shape.top - buffer,
                    "right": shape.left + shape.width + buffer,
                    "bottom": shape.top + shape.height + buffer
                })
    
    # Very small image size - max 15% of slide area
    max_area = slide_width * slide_height * 0.15
    aspect_ratio = image_data["width"] / image_data["height"] if image_data["height"] > 0 else 1.0
    
    # Calculate conservative size - convert to int to avoid float errors
    if aspect_ratio > 1:
        max_width = int(min(slide_width * 0.25, Inches(2.5)))
        max_height = int(max_width / aspect_ratio)
    else:
        max_height = int(min(slide_height * 0.25, Inches(2)))
        max_width = int(max_height * aspect_ratio)
    
    # Try corner positions that are least likely to interfere
    corner_positions = [
        # Top right corner
        {
            "left": int(slide_width - max_width - (slide_width * 0.02)),
            "top": int(slide_height * 0.02),
            "width": max_width,
            "height": max_height
        },
        # Bottom right corner
        {
            "left": int(slide_width - max_width - (slide_width * 0.02)),
            "top": int(slide_height - max_height - (slide_height * 0.02)),
            "width": max_width,
            "height": max_height
        },
        # Bottom left corner
        {
            "left": int(slide_width * 0.02),
            "top": int(slide_height - max_height - (slide_height * 0.02)),
            "width": max_width,
            "height": max_height
        }
    ]
    
    # Check each position for conflicts
    for position in corner_positions:
        conflicts = False
        for text_area in text_areas:
            if _areas_overlap(position, text_area):
                conflicts = True
                break
        
        if not conflicts:
            return position
    
    # If all corners conflict, don't add image
    return None

def _find_image_placeholder(slide):
    """Find an image placeholder in the slide."""
    for shape in slide.shapes:
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            # Common image placeholder types:
            # 18 = PP_PLACEHOLDER.PICTURE
            # 19 = PP_PLACEHOLDER.CLIP_ART  
            # 20 = PP_PLACEHOLDER.MEDIA_CLIP
            if placeholder_type in [18, 19, 20]:
                return shape
    
    # Also check for content placeholders that can hold images
    for shape in slide.shapes:
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            # 7 = PP_PLACEHOLDER.OBJECT (can be used for images)
            if placeholder_type == 7:
                return shape
    
    return None

def _replace_placeholder_with_image(slide, placeholder, image_path):
    """Replace a placeholder with an image while preserving position and size."""
    # Get placeholder properties
    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height
    
    # Remove the placeholder
    placeholder_element = placeholder.element
    placeholder_element.getparent().remove(placeholder_element)
    
    # Add image in the same position
    slide.shapes.add_picture(image_path, left, top, width, height)

def _calculate_image_placement(slide, image_data, layout_hint):
    """Calculate optimal image placement based on slide content and layout.

    This now uses the actual template slide dimensions rather than assuming 16:9.
    """
    # Get real slide dimensions from the slide part (already EMU units)
    try:
        slide_width = slide.part.slide_width
        slide_height = slide.part.slide_height
    except Exception:
        # Fallback to 16:9 default (10x7.5 inches in EMU)
        slide_width = Inches(10)
        slide_height = Inches(7.5)
    
    # Find available space by analyzing existing shapes and their types
    text_areas = []
    available_areas = []
    
    for shape in slide.shapes:
        if hasattr(shape, 'left') and hasattr(shape, 'top'):
            shape_info = {
                "left": shape.left,
                "top": shape.top,
                "right": shape.left + shape.width,
                "bottom": shape.top + shape.height,
                "is_text": hasattr(shape, 'text_frame'),
                "is_placeholder": shape.is_placeholder
            }
            
            if shape_info["is_text"]:
                text_areas.append(shape_info)
    
    # Calculate image size maintaining aspect ratio
    original_width = image_data["width"]
    original_height = image_data["height"]
    aspect_ratio = original_width / original_height if original_height > 0 else 1.0
    
    # Determine placement strategy based on layout hint and content
    placement_strategies = _get_placement_strategies(layout_hint, text_areas, aspect_ratio)
    
    # Try each placement strategy
    for strategy in placement_strategies:
        max_width = strategy["max_width"]
        max_height = strategy["max_height"]
        
        # Calculate final dimensions maintaining aspect ratio
        if aspect_ratio > max_width / max_height:
            final_width = max_width
            final_height = max_width / aspect_ratio
        else:
            final_height = max_height
            final_width = max_height * aspect_ratio
        
        # Check if this position conflicts with existing content
        conflicts = False
        test_area = {
            "left": strategy["left"],
            "top": strategy["top"],
            "right": strategy["left"] + final_width,
            "bottom": strategy["top"] + final_height
        }
        
        for text_area in text_areas:
            if _areas_overlap(test_area, text_area):
                conflicts = True
                break
        
        if not conflicts:
            return {
                "left": strategy["left"],
                "top": strategy["top"],
                "width": final_width,
                "height": final_height
            }
    
    # If no strategy works, use a safe fallback position
    return _get_fallback_placement(aspect_ratio, slide_width, slide_height)

def _get_placement_strategies(layout_hint, text_areas, aspect_ratio):
    """Get ordered list of placement strategies based on layout and content."""
    strategies = []
    
    # Analyze text areas to find the best placement
    if text_areas:
        # Find the main content area (usually the largest text area)
        main_text_area = max(text_areas, key=lambda x: (x["right"] - x["left"]) * (x["bottom"] - x["top"]))
        
        # Strategy 1: Right side of main text
        if main_text_area["right"] < Inches(6):
            strategies.append({
                "left": main_text_area["right"] + Inches(0.2),
                "top": main_text_area["top"],
                "max_width": Inches(9) - main_text_area["right"],
                "max_height": Inches(3),
                "name": "right_of_text"
            })
        
        # Strategy 2: Below main text
        if main_text_area["bottom"] < Inches(5):
            strategies.append({
                "left": main_text_area["left"],
                "top": main_text_area["bottom"] + Inches(0.2),
                "max_width": Inches(4),
                "max_height": Inches(6) - main_text_area["bottom"],
                "name": "below_text"
            })
    
    # Strategy 3: Layout hint-based positioning
    if "image-left" in layout_hint.lower():
        strategies.insert(0, {
            "left": Inches(0.5),
            "top": Inches(2.0),
            "max_width": Inches(4.0),
            "max_height": Inches(3.5),
            "name": "left_side"
        })
    elif "image-right" in layout_hint.lower():
        strategies.insert(0, {
            "left": Inches(5.5),
            "top": Inches(2.0),
            "max_width": Inches(4.0),
            "max_height": Inches(3.5),
            "name": "right_side"
        })
    elif "banner" in layout_hint.lower() or "section" in layout_hint.lower():
        strategies.insert(0, {
            "left": Inches(1.0),
            "top": Inches(1.5),
            "max_width": Inches(8.0),
            "max_height": Inches(2.0),
            "name": "banner"
        })
    
    # Strategy 4: Standard positions (fallbacks)
    standard_positions = [
        {
            "left": Inches(6.0),
            "top": Inches(2.5),
            "max_width": Inches(3.5),
            "max_height": Inches(2.5),
            "name": "top_right"
        },
        {
            "left": Inches(0.5),
            "top": Inches(4.0),
            "max_width": Inches(3.5),
            "max_height": Inches(2.5),
            "name": "bottom_left"
        },
        {
            "left": Inches(3.0),
            "top": Inches(4.5),
            "max_width": Inches(4.0),
            "max_height": Inches(2.0),
            "name": "bottom_center"
        }
    ]
    
    strategies.extend(standard_positions)
    return strategies

def _should_use_image_placeholder(slide, image_pool):
    """Determine if we should use an image placeholder for this slide."""
    if not image_pool:
        return False
    
    # Always prefer image placeholders when they exist
    for shape in slide.shapes:
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            if placeholder_type in [18, 19, 20]:  # Image placeholder types
                return True
    return False

def _get_slide_area(slide):
    """Get the total area of the slide in square inches."""
    try:
        width = slide.part.slide_width / 914400  # Convert EMU to inches
        height = slide.part.slide_height / 914400
        return width * height
    except Exception:
        return 75  # Default 10x7.5 inches

def _areas_overlap(area1, area2):
    """Check if two rectangular areas overlap."""
    # Convert all values to integers to avoid float type errors
    try:
        a1_left = int(area1["left"])
        a1_right = int(area1.get("right", area1["left"] + area1.get("width", 0)))
        a1_top = int(area1["top"])
        a1_bottom = int(area1.get("bottom", area1["top"] + area1.get("height", 0)))
        
        a2_left = int(area2["left"])
        a2_right = int(area2.get("right", area2["left"] + area2.get("width", 0)))
        a2_top = int(area2["top"])
        a2_bottom = int(area2.get("bottom", area2["top"] + area2.get("height", 0)))
        
        return not (a1_right <= a2_left or 
                    a1_left >= a2_right or
                    a1_bottom <= a2_top or 
                    a1_top >= a2_bottom)
    except (TypeError, ValueError):
        # Fallback to safe non-overlap if conversion fails
        return False

def _get_fallback_placement(aspect_ratio, slide_width, slide_height):
    """Get a safe fallback placement using actual slide dimensions.

    Places the image in the lower-right quadrant with size ~30% of slide width.
    """
    try:
        # Use proportions of slide size (slide_width/height are EMU or Length objects)
        max_width = int(slide_width * 0.3)
        max_height = int(slide_height * 0.25)
    except Exception:
        max_width = Inches(3.0)
        max_height = Inches(2.0)

    if aspect_ratio > max_width / max_height:
        final_width = max_width
        final_height = int(max_width / aspect_ratio)
    else:
        final_height = max_height
        final_width = int(max_height * aspect_ratio)

    return {
        "left": int(slide_width - final_width - (slide_width * 0.05)),  # 5% margin
        "top": int(slide_height - final_height - (slide_height * 0.05)),
        "width": final_width,
        "height": final_height
    }

def _parse_markdown_formatting(text):
    """Parse text with Markdown formatting and return list of (text, is_bold, is_italic) tuples."""
    parts = []
    current_pos = 0
    
    # Pattern for **bold** and *italic*
    bold_pattern = r'\*\*(.*?)\*\*'
    italic_pattern = r'\*(.*?)\*'
    
    # Find all bold matches first
    bold_matches = list(re.finditer(bold_pattern, text))
    italic_matches = list(re.finditer(italic_pattern, text))
    
    # Remove italic matches that are inside bold matches (to avoid double processing **)
    filtered_italic_matches = []
    for italic_match in italic_matches:
        is_inside_bold = False
        for bold_match in bold_matches:
            if bold_match.start() <= italic_match.start() and italic_match.end() <= bold_match.end():
                is_inside_bold = True
                break
        if not is_inside_bold:
            filtered_italic_matches.append(italic_match)
    
    # Combine and sort all matches
    all_matches = []
    for match in bold_matches:
        all_matches.append((match.start(), match.end(), match.group(1), True, False))
    for match in filtered_italic_matches:
        all_matches.append((match.start(), match.end(), match.group(1), False, True))
    
    all_matches.sort(key=lambda x: x[0])
    
    # Build parts list
    for start, end, matched_text, is_bold, is_italic in all_matches:
        # Add text before this match
        if current_pos < start:
            plain_text = text[current_pos:start]
            if plain_text:
                parts.append((plain_text, False, False))
        
        # Add the formatted text
        parts.append((matched_text, is_bold, is_italic))
        current_pos = end
    
    # Add remaining text
    if current_pos < len(text):
        remaining_text = text[current_pos:]
        if remaining_text:
            parts.append((remaining_text, False, False))
    
    # If no formatting was found, return the whole text as plain
    if not parts:
        parts = [(text, False, False)]
    
    return parts

def _apply_font_style(shape, font_family):
    """Apply font styling to a text shape."""
    if not hasattr(shape, 'text_frame'):
        return
    
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font:
            paragraph.font.name = font_family


