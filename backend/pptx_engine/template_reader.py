from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER

def analyze_template(path: str):
    """Analyze a PowerPoint template to extract reusable style and assets."""
    prs = Presentation(path)

    # Collect layouts with detailed placeholder information
    layouts = {}
    layout_info = {}
    
    for i, layout in enumerate(prs.slide_layouts):
        layouts[i] = layout
        
        # Analyze placeholders in each layout
        placeholders = []
        for shape in layout.shapes:
            if shape.is_placeholder:
                placeholder_info = {
                    "type": shape.placeholder_format.type,
                    "idx": shape.placeholder_format.idx,
                    "name": getattr(shape, 'name', f'placeholder_{shape.placeholder_format.idx}'),
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                    "has_text_frame": hasattr(shape, 'text_frame'),
                    "is_image_placeholder": shape.placeholder_format.type in [18, 19, 20]  # Picture placeholders
                }
                
                # Extract font information from text placeholders
                if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                    font_info = _extract_font_info(shape.text_frame.paragraphs[0])
                    placeholder_info["font"] = font_info
                
                placeholders.append(placeholder_info)
        
        layout_info[i] = {
            "name": layout.name,
            "placeholders": placeholders
        }
    
    # Enhanced theme extraction
    theme = {
        "font_family": _infer_font(prs),
        "accent_colors": _infer_colors(prs),
    "layout_info": layout_info,
    # Store slide dimensions (English Metric Units from pptx)
    "slide_width_emu": prs.slide_width,
    "slide_height_emu": prs.slide_height,
    "slide_width_inches": float(prs.slide_width) / 914400.0,
    "slide_height_inches": float(prs.slide_height) / 914400.0,
    "aspect_ratio": round((float(prs.slide_width) / float(prs.slide_height)), 4) if prs.slide_height else 1.0
    }

    # Harvest images from all slides in template with better metadata
    reusable_images = []
    for s_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_info = {
                        "slide_index": s_idx,
                        "name": getattr(shape, "name", f"pic_{s_idx}"),
                        "blob": shape.image.blob,
                        "width": shape.width, 
                        "height": shape.height,
                        "left": shape.left,
                        "top": shape.top,
                        "aspect_ratio": shape.width / shape.height if shape.height > 0 else 1.0
                    }
                    reusable_images.append(image_info)
                except Exception:
                    continue
    
    return {
        "layouts": layouts,
        "theme": theme,
        "images": reusable_images,
        "prs": prs
    }

def _extract_font_info(paragraph):
    """Extract detailed font information from a paragraph."""
    font_info = {
        "name": "Calibri",
        "size": 18,
        "bold": False,
        "italic": False,
        "color": None
    }
    
    try:
        if paragraph.font:
            if paragraph.font.name:
                font_info["name"] = paragraph.font.name
            if paragraph.font.size:
                font_info["size"] = paragraph.font.size.pt
            if paragraph.font.bold is not None:
                font_info["bold"] = paragraph.font.bold
            if paragraph.font.italic is not None:
                font_info["italic"] = paragraph.font.italic
            # Note: color extraction is complex due to theme colors
    except Exception:
        pass
    
    return font_info

def _infer_font(prs):
    """Infer the primary font from the template."""
    try:
        # Look at title placeholder on first layout as a heuristic
        layout = prs.slide_layouts[0]
        if hasattr(layout.shapes, 'title') and layout.shapes.title:
            text_frame = layout.shapes.title.text_frame
            if text_frame.paragraphs:
                font = text_frame.paragraphs[0].font
                if font and font.name:
                    return font.name
    except Exception:
        pass
    return "Calibri"

def _infer_colors(prs):
    """Infer accent colors from the template."""
    colors = []
    try:
        # Look at fills from shapes in first few layouts
        for layout in prs.slide_layouts[:3]:
            for shape in layout.shapes:
                if hasattr(shape, "fill") and shape.fill:
                    try:
                        if hasattr(shape.fill, "fore_color") and shape.fill.fore_color:
                            if hasattr(shape.fill.fore_color, "rgb"):
                                colors.append(str(shape.fill.fore_color.rgb))
                    except Exception:
                        continue
                if len(colors) >= 6:
                    break
            if len(colors) >= 6:
                break
    except Exception:
        pass
    return colors[:6]
