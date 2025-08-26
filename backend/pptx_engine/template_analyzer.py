"""
Template Analyzer - Extracts template structure for LLM guidance
"""
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
import json

def analyze_template_structure(template_path: str) -> dict:
    """Analyze template and extract layout structures for LLM guidance."""
    prs = Presentation(template_path)
    
    template_structure = {
        "slide_dimensions": {
            "width_inches": float(prs.slide_width) / 914400.0,
            "height_inches": float(prs.slide_height) / 914400.0,
            "aspect_ratio": round((float(prs.slide_width) / float(prs.slide_height)), 4)
        },
        "layouts": []
    }
    
    # Analyze each layout
    for layout_idx, layout in enumerate(prs.slide_layouts):
        layout_info = {
            "layout_index": layout_idx,
            "layout_name": layout.name,
            "placeholders": [],
            "structure_description": "",
            "recommended_content_types": []
        }
        
        # Analyze placeholders in this layout
        placeholder_summary = []
        has_title = False
        has_content = False
        has_image = False
        has_subtitle = False
        
        for shape in layout.shapes:
            if shape.is_placeholder:
                placeholder_type = shape.placeholder_format.type
                placeholder_info = _analyze_placeholder(shape, placeholder_type)
                layout_info["placeholders"].append(placeholder_info)
                
                # Track what this layout has
                if placeholder_type == 1:  # Title
                    has_title = True
                    placeholder_summary.append("title")
                elif placeholder_type == 2:  # Body/Content
                    has_content = True
                    placeholder_summary.append("content")
                elif placeholder_type == 3:  # Subtitle
                    has_subtitle = True
                    placeholder_summary.append("subtitle")
                elif placeholder_type in [18, 19, 20]:  # Image placeholders
                    has_image = True
                    placeholder_summary.append("image")
                elif placeholder_type == 7:  # Object (can be image)
                    placeholder_summary.append("object/image")
        
        # Generate structure description for LLM
        layout_info["structure_description"] = _generate_layout_description(
            layout.name, placeholder_summary, has_title, has_content, has_image, has_subtitle
        )
        
        # Recommend content types
        layout_info["recommended_content_types"] = _get_recommended_content_types(
            layout.name, has_title, has_content, has_image, has_subtitle
        )
        
        template_structure["layouts"].append(layout_info)
    
    return template_structure

def _analyze_placeholder(shape, placeholder_type):
    """Analyze individual placeholder properties."""
    # Convert position and size to relative percentages for LLM understanding
    try:
        slide_width = shape.part.slide_width if hasattr(shape, 'part') else 9144000  # 10 inches in EMU
        slide_height = shape.part.slide_height if hasattr(shape, 'part') else 6858000  # 7.5 inches in EMU
    except:
        slide_width = 9144000
        slide_height = 6858000
    
    placeholder_info = {
        "type": placeholder_type,
        "type_name": _get_placeholder_type_name(placeholder_type),
        "position": {
            "left_percent": round((shape.left / slide_width) * 100, 1),
            "top_percent": round((shape.top / slide_height) * 100, 1),
            "width_percent": round((shape.width / slide_width) * 100, 1),
            "height_percent": round((shape.height / slide_height) * 100, 1)
        },
        "size_category": _categorize_placeholder_size(shape.width, shape.height, slide_width, slide_height),
        "likely_content": _suggest_placeholder_content(placeholder_type)
    }
    
    return placeholder_info

def _get_placeholder_type_name(placeholder_type):
    """Get human-readable placeholder type name."""
    type_names = {
        1: "title",
        2: "body_content", 
        3: "subtitle",
        7: "object_or_image",
        13: "text",
        14: "content",
        18: "picture",
        19: "clip_art",
        20: "media"
    }
    return type_names.get(placeholder_type, f"type_{placeholder_type}")

def _categorize_placeholder_size(width, height, slide_width, slide_height):
    """Categorize placeholder size for LLM understanding."""
    width_percent = (width / slide_width) * 100
    height_percent = (height / slide_height) * 100
    area_percent = width_percent * height_percent / 100
    
    if area_percent > 50:
        return "large"
    elif area_percent > 25:
        return "medium" 
    elif area_percent > 10:
        return "small"
    else:
        return "tiny"

def _suggest_placeholder_content(placeholder_type):
    """Suggest what content should go in this placeholder."""
    suggestions = {
        1: "Main slide title or heading",
        2: "Bullet points, main content, or body text",
        3: "Subtitle, tagline, or secondary heading", 
        7: "Images, charts, or other objects",
        13: "Additional text or supporting content",
        14: "Flexible content area",
        18: "Pictures or photographs",
        19: "Clip art or illustrations", 
        20: "Media content (videos, audio)"
    }
    return suggestions.get(placeholder_type, "Content area")

def _generate_layout_description(layout_name, placeholder_summary, has_title, has_content, has_image, has_subtitle):
    """Generate human-readable layout description for LLM."""
    description_parts = []
    
    # Start with layout name if meaningful
    if layout_name and layout_name not in ["Layout", "Slide Layout"]:
        description_parts.append(f"'{layout_name}' layout")
    
    # Describe structure
    if has_title and has_content and has_image:
        description_parts.append("title + content + image areas")
    elif has_title and has_content:
        description_parts.append("title + content areas")
    elif has_title and has_image:
        description_parts.append("title + image areas")
    elif has_title and has_subtitle:
        description_parts.append("title + subtitle areas")
    elif has_title:
        description_parts.append("title-focused layout")
    elif has_content:
        description_parts.append("content-focused layout")
    
    # Add placeholder summary
    if placeholder_summary:
        unique_placeholders = list(set(placeholder_summary))
        description_parts.append(f"Contains: {', '.join(unique_placeholders)}")
    
    return " - ".join(description_parts) if description_parts else "Basic layout"

def _get_recommended_content_types(layout_name, has_title, has_content, has_image, has_subtitle):
    """Get recommended content types for this layout."""
    recommendations = []
    
    layout_lower = layout_name.lower() if layout_name else ""
    
    # Based on layout name patterns
    if "title" in layout_lower and "content" in layout_lower:
        recommendations.extend(["presentation_slides", "bullet_points", "general_content"])
    elif "title" in layout_lower:
        recommendations.extend(["section_headers", "chapter_titles", "intro_slides"])
    elif "section" in layout_lower:
        recommendations.extend(["section_breaks", "chapter_dividers"])
    elif "comparison" in layout_lower or "two" in layout_lower:
        recommendations.extend(["comparisons", "before_after", "pros_cons"])
    elif "picture" in layout_lower or "image" in layout_lower:
        recommendations.extend(["image_showcases", "photo_galleries", "visual_content"])
    
    # Based on structure
    if has_title and has_content and has_image:
        recommendations.extend(["standard_slides", "mixed_content"])
    elif has_title and has_content:
        recommendations.extend(["text_heavy_slides", "bullet_lists"])
    elif has_image:
        recommendations.extend(["visual_slides", "image_content"])
    
    return list(set(recommendations)) if recommendations else ["general_purpose"]

def generate_llm_template_prompt(template_structure: dict, user_content: str) -> str:
    """Generate a detailed prompt for LLM with template structure information."""
    
    prompt = f"""You are creating a PowerPoint presentation using a specific template. Here's the template structure:

TEMPLATE INFORMATION:
- Slide dimensions: {template_structure['slide_dimensions']['width_inches']:.1f}" × {template_structure['slide_dimensions']['height_inches']:.1f}" (aspect ratio: {template_structure['slide_dimensions']['aspect_ratio']})

AVAILABLE LAYOUTS:
"""
    
    for layout in template_structure['layouts']:
        prompt += f"\nLayout {layout['layout_index']}: {layout['structure_description']}\n"
        prompt += f"  Recommended for: {', '.join(layout['recommended_content_types'])}\n"
        prompt += f"  Available placeholders:\n"
        
        for placeholder in layout['placeholders']:
            prompt += f"    - {placeholder['type_name']}: {placeholder['likely_content']} "
            prompt += f"({placeholder['size_category']} size, {placeholder['position']['width_percent']:.0f}%×{placeholder['position']['height_percent']:.0f}%)\n"
    
    prompt += f"""
USER CONTENT TO ORGANIZE:
{user_content}

INSTRUCTIONS:
Please create a presentation plan by:
1. Analyzing the user content and determining how many slides are needed
2. For each slide, choose the most appropriate layout from the available options
3. Specify exactly which content goes in which placeholder
4. Indicate if images should be used and where

Respond with a JSON structure like this:
{{
  "slides": [
    {{
      "layout_index": 0,
      "layout_reasoning": "Why this layout was chosen",
      "placeholders": {{
        "title": "Exact text for title placeholder",
        "body_content": ["Bullet point 1", "Bullet point 2"],
        "subtitle": "Subtitle text if needed",
        "picture": {{"use_image": true, "description": "What kind of image"}},
        "object_or_image": {{"use_image": false}}
      }}
    }}
  ]
}}

Make sure to:
- Choose layouts that best fit the content type
- Don't overfill placeholders - keep text concise
- Use image placeholders strategically
- Ensure content flows logically across slides
"""
    
    return prompt

if __name__ == "__main__":
    # Test the analyzer
    template_path = "path/to/template.pptx"
    structure = analyze_template_structure(template_path)
    print(json.dumps(structure, indent=2))
