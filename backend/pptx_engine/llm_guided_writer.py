"""
LLM-Guided Slide Writer - Uses template analysis to guide content placement
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import tempfile
import os
import json
from io import BytesIO
from .template_analyzer import analyze_template_structure, generate_llm_template_prompt

def build_presentation_with_llm_guidance(plan, style_ctx, template_path: str, llm_client) -> bytes:
    """Build a PowerPoint presentation using LLM guidance based on template structure."""
    
    # Step 1: Analyze the template structure
    template_structure = analyze_template_structure(template_path)
    
    # Step 2: Get user content from the plan
    user_content = _extract_user_content_from_plan(plan)
    
    # Step 3: Generate LLM prompt with template structure
    llm_prompt = generate_llm_template_prompt(template_structure, user_content)
    
    # Step 4: Get LLM guidance on content placement
    try:
        llm_response = llm_client.get_completion(llm_prompt)
        llm_plan = json.loads(llm_response)
    except Exception as e:
        print(f"LLM guidance failed, falling back to original plan: {e}")
        llm_plan = _convert_original_plan_to_llm_format(plan)
    
    # Step 5: Build presentation following LLM guidance exactly
    return _build_presentation_from_llm_plan(llm_plan, template_structure, style_ctx, template_path)

def _extract_user_content_from_plan(plan):
    """Extract all user content from the original plan."""
    content_parts = []
    
    # Add any top-level information
    if "title" in plan:
        content_parts.append(f"Presentation Title: {plan['title']}")
    
    if "description" in plan:
        content_parts.append(f"Description: {plan['description']}")
    
    # Extract content from all slides
    for i, slide in enumerate(plan.get("slides", [])):
        content_parts.append(f"\nSlide {i+1} Content:")
        
        if "title" in slide:
            content_parts.append(f"Title: {slide['title']}")
        
        if "bullets" in slide and slide["bullets"]:
            content_parts.append("Key Points:")
            for bullet in slide["bullets"]:
                content_parts.append(f"- {bullet}")
        
        if "notes" in slide:
            content_parts.append(f"Speaker Notes: {slide['notes']}")
    
    return "\n".join(content_parts)

def _convert_original_plan_to_llm_format(plan):
    """Convert original plan to LLM format as fallback."""
    llm_plan = {"slides": []}
    
    for slide in plan.get("slides", []):
        llm_slide = {
            "layout_index": 1,  # Default to content layout
            "layout_reasoning": "Fallback to standard content layout",
            "placeholders": {}
        }
        
        if "title" in slide:
            llm_slide["placeholders"]["title"] = slide["title"]
        
        if "bullets" in slide:
            llm_slide["placeholders"]["body_content"] = slide["bullets"]
        
        llm_plan["slides"].append(llm_slide)
    
    return llm_plan

def _build_presentation_from_llm_plan(llm_plan, template_structure, style_ctx, template_path):
    """Build presentation following LLM guidance exactly."""
    # Create new presentation from template
    prs = Presentation(template_path)
    
    # Clear existing slides
    slide_idxs = list(range(len(prs.slides) - 1, -1, -1))
    for i in slide_idxs:
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    
    # Prepare image pool
    image_pool = style_ctx.get("images", [])
    image_usage_count = 0
    
    # Create slides according to LLM plan
    for slide_plan in llm_plan.get("slides", []):
        # Get the specified layout
        layout_index = slide_plan.get("layout_index", 1)
        
        # Ensure layout index is valid
        if layout_index >= len(prs.slide_layouts):
            layout_index = 1  # Fallback to content layout
        
        layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(layout)
        
        # Fill placeholders exactly as specified by LLM
        placeholders_plan = slide_plan.get("placeholders", {})
        _fill_slide_according_to_llm_plan(slide, placeholders_plan, template_structure, layout_index, image_pool, image_usage_count)
        
        # Update image usage counter if images were used
        if any(p.get("use_image", False) for p in placeholders_plan.values() if isinstance(p, dict)):
            image_usage_count += 1
    
    # Save to bytes
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()

def _fill_slide_according_to_llm_plan(slide, placeholders_plan, template_structure, layout_index, image_pool, image_usage_count):
    """Fill slide placeholders exactly according to LLM plan."""
    
    # Get layout structure information
    layout_info = template_structure["layouts"][layout_index] if layout_index < len(template_structure["layouts"]) else {}
    
    # Map LLM placeholder names to actual placeholder types
    placeholder_type_mapping = {
        "title": 1,
        "body_content": 2,
        "subtitle": 3,
        "object_or_image": 7,
        "text": 13,
        "content": 14,
        "picture": 18,
        "clip_art": 19,
        "media": 20
    }
    
    # Fill each placeholder as specified
    for shape in slide.shapes:
        if not (hasattr(shape, 'is_placeholder') and shape.is_placeholder):
            continue
        
        placeholder_type = shape.placeholder_format.type
        
        # Find which LLM placeholder corresponds to this shape
        llm_placeholder_name = None
        for llm_name, type_num in placeholder_type_mapping.items():
            if type_num == placeholder_type and llm_name in placeholders_plan:
                llm_placeholder_name = llm_name
                break
        
        if not llm_placeholder_name:
            continue
        
        placeholder_content = placeholders_plan[llm_placeholder_name]
        
        # Handle different content types
        if isinstance(placeholder_content, str):
            # Simple text content
            _fill_text_placeholder_precisely(shape, placeholder_content, placeholder_type)
        
        elif isinstance(placeholder_content, list):
            # Bullet points or list content
            _fill_bullet_placeholder_precisely(shape, placeholder_content)
        
        elif isinstance(placeholder_content, dict):
            # Image or complex content
            if placeholder_content.get("use_image", False) and image_pool:
                _add_image_to_placeholder(shape, slide, image_pool, image_usage_count, placeholder_content)
            elif "text" in placeholder_content:
                _fill_text_placeholder_precisely(shape, placeholder_content["text"], placeholder_type)

def _fill_text_placeholder_precisely(shape, text, placeholder_type):
    """Fill text placeholder with precise formatting."""
    if not hasattr(shape, 'text_frame') or not text:
        return
    
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    # Add text
    paragraph = text_frame.paragraphs[0]
    paragraph.text = str(text)
    
    # Apply appropriate formatting based on placeholder type
    if placeholder_type == 1:  # Title
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
    elif placeholder_type == 3:  # Subtitle
        paragraph.font.size = Pt(20)
        paragraph.font.bold = False
    else:  # Body text
        paragraph.font.size = Pt(16)
        paragraph.font.bold = False

def _fill_bullet_placeholder_precisely(shape, bullet_list):
    """Fill placeholder with bullet points precisely."""
    if not hasattr(shape, 'text_frame') or not bullet_list:
        return
    
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    for i, bullet_text in enumerate(bullet_list):
        if i == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        
        paragraph.text = str(bullet_text)
        paragraph.level = 0
        paragraph.font.size = Pt(16)

def _add_image_to_placeholder(shape, slide, image_pool, image_usage_count, image_spec):
    """Add image to placeholder following LLM specification."""
    if not image_pool:
        return
    
    # Select image based on usage pattern
    image_idx = image_usage_count % len(image_pool)
    image_data = image_pool[image_idx]
    
    try:
        # Create temporary file for the image
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
            tmp_file.write(image_data["blob"])
            tmp_path = tmp_file.name
        
        # Get placeholder properties
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        
        # Remove the placeholder
        placeholder_element = shape.element
        placeholder_element.getparent().remove(placeholder_element)
        
        # Add image in the same position
        slide.shapes.add_picture(tmp_path, left, top, width, height)
        
    except Exception as e:
        print(f"Failed to add image to placeholder: {e}")
    finally:
        try:
            os.remove(tmp_path)
        except:
            pass

def create_llm_client():
    """Create LLM client - placeholder for actual implementation."""
    class MockLLMClient:
        def get_completion(self, prompt):
            # This would be replaced with actual LLM API call
            return '{"slides": [{"layout_index": 1, "layout_reasoning": "Standard content layout", "placeholders": {"title": "Sample Title", "body_content": ["Point 1", "Point 2"]}}]}'
    
    return MockLLMClient()
