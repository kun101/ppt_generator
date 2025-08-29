"""Built-in default PowerPoint template generator.

Provides a lightweight branded-looking template so users can try the
service without supplying their own .pptx/.potx file.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def get_default_template_bytes() -> bytes:
    """Create a simple default template and return raw pptx bytes.

    The template includes:
    - Title slide
    - Title + Content slide layout
    - Section header slide
    """
    prs = Presentation()

    # Ensure 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Default Presentation Template"
    subtitle.text = "Powered by AI Layout Guidance"
    _style_title_shape(title)
    _style_subtitle_shape(subtitle)

    # Content slide
    content_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(content_layout)
    title2 = slide2.shapes.title
    body2 = slide2.placeholders[1]
    title2.text = "Example Content Slide"
    tf = body2.text_frame
    tf.clear()
    for bullet in [
        "Reusable default theme",
        "Safe text fitting & no overflow",
        "Deterministic placeholder mapping",
        "Smart template analysis"
    ]:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = bullet
        p.level = 0
    _style_title_shape(title2)

    # Section header (if available)
    if len(prs.slide_layouts) > 2:
        section_layout = prs.slide_layouts[2]
        section_slide = prs.slides.add_slide(section_layout)
        if section_slide.shapes.title:
            section_slide.shapes.title.text = "Section Header"
            _style_title_shape(section_slide.shapes.title)

    bio = _save_to_bytes(prs)
    return bio


def _style_title_shape(shape):
    if not getattr(shape, 'text_frame', None):
        return
    p = shape.text_frame.paragraphs[0]
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(33, 150, 243)  # Blue accent


def _style_subtitle_shape(shape):
    if not getattr(shape, 'text_frame', None):
        return
    p = shape.text_frame.paragraphs[0]
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(200, 200, 200)


def _save_to_bytes(prs: Presentation) -> bytes:
    from io import BytesIO
    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()


def get_default_template_metadata():
    """Return a lightweight preview/metadata structure for UI preview."""
    return {
        "name": "Built-in Default Template",
        "description": "Clean 16:9 theme with title, content, and section layouts.",
        "layouts": [
            {"index": 0, "type": "title", "placeholders": ["title", "subtitle"]},
            {"index": 1, "type": "title+content", "placeholders": ["title", "body"]},
            {"index": 2, "type": "section", "placeholders": ["title"]},
        ],
        "colors": ["#2196F3", "#FFFFFF", "#C8C8C8"],
        "aspect_ratio": "16:9"
    }


__all__ = [
    "get_default_template_bytes",
    "get_default_template_metadata"
]
