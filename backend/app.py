from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse, JSONResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from backend.utils.security import scrub_logs
from backend.pptx_engine.template_reader import analyze_template
from backend.pptx_engine.template_analyzer import analyze_template_structure, generate_llm_template_prompt
from backend.pptx_engine.slide_planner import make_slide_plan
from backend.pptx_engine.slide_writer import build_presentation
from backend.pptx_engine.default_template import (
    get_default_template_bytes,
    get_default_template_metadata,
)
from backend.pptx_engine.llm_guided_writer import build_presentation_with_llm_guidance
from backend.llm import get_provider
from tempfile import NamedTemporaryFile
import io
import os
import logging
import json
from pathlib import Path

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="Text→PPTX", description="Convert text to PowerPoint presentations")

# Log startup
logger.info("Starting Text-to-PowerPoint Generator")
logger.info(f"Python path: {os.getcwd()}")
logger.info(f"PORT environment variable: {os.environ.get('PORT', 'Not set')}")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], 
    allow_credentials=False,
    allow_methods=["POST", "GET", "OPTIONS"], 
    allow_headers=["*"]
)

@app.get("/api/health")
def health():
    logger.info("Health check endpoint called")
    return {"ok": True, "service": "text-to-pptx", "status": "healthy"}

@app.get("/api/info")
async def api_info():
    """API info endpoint."""
    return {
        "message": "Text-to-PowerPoint Generator API",
        "status": "running",
        "health": "/api/health",
        "default_template": "/api/template/default",
        "default_template_meta": "/api/template/default/meta"
    }

@app.get("/api/debug/routes")
def debug_routes():
    """Debug endpoint to show all available routes."""
    routes = []
    for route in app.routes:
        routes.append({
            "path": getattr(route, 'path', 'N/A'),
            "methods": getattr(route, 'methods', 'N/A'),
            "name": getattr(route, 'name', 'N/A')
        })
    return {"routes": routes}


@app.get("/api/template/default/meta")
def default_template_meta():
    """Metadata / preview JSON for built-in default template."""
    return get_default_template_metadata()


@app.get("/api/template/default")
def default_template_download():
    """Download the built-in default template (.pptx)."""
    data = get_default_template_bytes()
    return StreamingResponse(
        io.BytesIO(data),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="default_template.pptx"'}
    )

@app.post("/api/generate")
async def generate_pptx(
    text: str = Form(...),
    guidance: str = Form(""),
    provider: str = Form(...),        # "openai" | "gemini"
    api_key: str = Form(...),
    template: UploadFile = File(...),
    use_template_guidance: bool = Form(True)  # Default now True - smart template analysis
):
    print(f"DEBUG: Received request with provider={provider}")  # Debug log
    scrub_logs()  # ensure no sensitive logs

    # Validate file type
    if not template.filename.lower().endswith(('.pptx', '.potx')):
        raise HTTPException(status_code=400, detail="Template must be a .pptx or .potx file")

    # Save template to a temp file
    with NamedTemporaryFile(delete=False, suffix=os.path.splitext(template.filename)[1]) as tf:
        tf.write(await template.read())
        template_path = tf.name

    try:
        # Always prefer template-guided unless explicitly disabled
        if use_template_guidance:
            return await _generate_with_template_guidance(text, guidance, template_path, provider, api_key)
        # Fallback explicit request for legacy mode
        return await _generate_standard(text, guidance, template_path, provider, api_key)
            
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))
    finally:
        try: 
            os.remove(template_path)
        except: 
            pass

async def _generate_with_template_guidance(text: str, guidance: str, template_path: str, provider: str, api_key: str):
    """Generate presentation using template structure analysis and LLM guidance."""
    
    # 1) Analyze template structure in detail
    template_structure = analyze_template_structure(template_path)
    
    # 2) Create LLM prompt with template structure information
    combined_content = f"{text}\n\nAdditional Guidance: {guidance}" if guidance else text
    llm_prompt = generate_llm_template_prompt(template_structure, combined_content)
    
    # 3) Get LLM guidance on exact content placement
    llm = get_provider(provider, api_key)
    try:
        # Get structured response from LLM
        llm_response = await llm.get_completion(llm_prompt)
        print(f"DEBUG: LLM template guidance received: {llm_response[:200]}...")
        
        # Parse LLM response with improved error handling
        try:
            llm_plan = json.loads(llm_response)
        except json.JSONDecodeError as json_error:
            print(f"DEBUG: JSON parsing failed: {json_error}")
            print(f"DEBUG: Raw response length: {len(llm_response)}")
            print(f"DEBUG: Response ends with: {llm_response[-100:] if len(llm_response) > 100 else llm_response}")
            raise Exception(f"Invalid JSON response: {json_error}")
        
    except Exception as e:
        print(f"DEBUG: LLM guidance failed, falling back: {e}")
        # Fallback to standard approach if LLM guidance fails
        return await _generate_standard(text, guidance, template_path, provider, api_key)
    
    # 4) Analyze template for assets (images, fonts, etc.)
    style_ctx = analyze_template(template_path)
    
    # 5) Build presentation following LLM guidance exactly
    pptx_bytes = _build_presentation_from_llm_guidance(llm_plan, template_structure, style_ctx, template_path)
    
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="template_guided.pptx"'}
    )

def _build_presentation_from_llm_guidance(llm_plan, template_structure, style_ctx, template_path):
    """Build presentation following LLM guidance exactly."""
    from pptx import Presentation
    from pptx.util import Pt
    import tempfile
    import os
    
    # Create new presentation from template
    prs = Presentation(template_path)
    
    # Clear existing slides
    slide_idxs = list(range(len(prs.slides) - 1, -1, -1))
    for i in slide_idxs:
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    
    # Prepare image pool
    image_pool = style_ctx.get("images", [])
    image_usage_count = 0
    
    # Create slides according to LLM plan
    for slide_plan in llm_plan.get("slides", []):
        # Get the specified layout
        layout_index = slide_plan.get("layout_index", 1)
        
        # Ensure layout index is valid
        if layout_index >= len(prs.slide_layouts):
            layout_index = 1  # Fallback to content layout
        
        layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(layout)
        
        # Fill placeholders exactly as specified by LLM
        placeholders_plan = slide_plan.get("placeholders", {})
        _fill_slide_according_to_llm_plan(slide, placeholders_plan, template_structure, layout_index, image_pool, image_usage_count)
        
        # Update image usage counter if images were used
        if any(p.get("use_image", False) for p in placeholders_plan.values() if isinstance(p, dict)):
            image_usage_count += 1
    
    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()

def _fill_slide_according_to_llm_plan(slide, placeholders_plan, template_structure, layout_index, image_pool, image_usage_count):
    """Fill slide placeholders exactly according to LLM plan."""
    
    # Map LLM placeholder names to actual placeholder types
    placeholder_type_mapping = {
        "title": 1,
        "body_content": 2,
        "subtitle": 3,
        "object_or_image": 7,
        "text": 13,
        "content": 14,
        "picture": 18,
        "clip_art": 19,
        "media": 20
    }
    
    # Fill each placeholder as specified
    for shape in slide.shapes:
        if not (hasattr(shape, 'is_placeholder') and shape.is_placeholder):
            continue
        
        placeholder_type = shape.placeholder_format.type
        
        # Find which LLM placeholder corresponds to this shape
        llm_placeholder_name = None
        for llm_name, type_num in placeholder_type_mapping.items():
            if type_num == placeholder_type and llm_name in placeholders_plan:
                llm_placeholder_name = llm_name
                break
        
        if not llm_placeholder_name:
            continue
        
        placeholder_content = placeholders_plan[llm_placeholder_name]
        
        # Handle different content types
        if isinstance(placeholder_content, str):
            # Simple text content
            _fill_text_placeholder_precisely(shape, placeholder_content, placeholder_type)
        
        elif isinstance(placeholder_content, list):
            # Bullet points or list content
            _fill_bullet_placeholder_precisely(shape, placeholder_content)
        
        elif isinstance(placeholder_content, dict):
            # Image or complex content
            if placeholder_content.get("use_image", False) and image_pool:
                _add_image_to_placeholder(shape, slide, image_pool, image_usage_count, placeholder_content)
            elif "text" in placeholder_content:
                _fill_text_placeholder_precisely(shape, placeholder_content["text"], placeholder_type)

def _fill_text_placeholder_precisely(shape, text, placeholder_type):
    """Fill text placeholder with precise formatting."""
    from pptx.util import Pt
    
    if not hasattr(shape, 'text_frame') or not text:
        return
    
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    # Add text
    paragraph = text_frame.paragraphs[0]
    paragraph.text = str(text)
    
    # Apply appropriate formatting based on placeholder type
    if placeholder_type == 1:  # Title
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
    elif placeholder_type == 3:  # Subtitle
        paragraph.font.size = Pt(20)
        paragraph.font.bold = False
    else:  # Body text
        paragraph.font.size = Pt(16)
        paragraph.font.bold = False

def _fill_bullet_placeholder_precisely(shape, bullet_list):
    """Fill placeholder with bullet points precisely."""
    from pptx.util import Pt
    
    if not hasattr(shape, 'text_frame') or not bullet_list:
        return
    
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    for i, bullet_text in enumerate(bullet_list):
        if i == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        
        paragraph.text = str(bullet_text)
        paragraph.level = 0
        paragraph.font.size = Pt(16)

def _add_image_to_placeholder(shape, slide, image_pool, image_usage_count, image_spec):
    """Add image to placeholder following LLM specification."""
    import tempfile
    import os
    
    if not image_pool:
        return
    
    # Select image based on usage pattern
    image_idx = image_usage_count % len(image_pool)
    image_data = image_pool[image_idx]
    
    try:
        # Create temporary file for the image
        with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_file:
            tmp_file.write(image_data["blob"])
            tmp_path = tmp_file.name
        
        # Get placeholder properties
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        
        # Remove the placeholder
        placeholder_element = shape.element
        placeholder_element.getparent().remove(placeholder_element)
        
        # Add image in the same position
        slide.shapes.add_picture(tmp_path, left, top, width, height)
        
    except Exception as e:
        print(f"Failed to add image to placeholder: {e}")
    finally:
        try:
            os.remove(tmp_path)
        except:
            pass

async def _generate_standard(text: str, guidance: str, template_path: str, provider: str, api_key: str):
    """Generate presentation using standard approach."""
    
    # 1) Read template style & assets
    style_ctx = analyze_template(template_path)

    # 2) Plan slides via LLM (with heuristic fallback)
    llm = get_provider(provider, api_key)
    slide_plan = await make_slide_plan(text, guidance, style_ctx, llm)

    # 3) Build PPTX
    pptx_bytes = build_presentation(slide_plan, style_ctx, template_path)

    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": 'attachment; filename="generated.pptx"'}
    )

# Setup paths for static files
current_dir = Path(__file__).parent.parent  # Go up from backend/ to root
frontend_dir = current_dir / "frontend"

logger.info(f"Current directory: {current_dir}")
logger.info(f"Frontend directory: {frontend_dir}")
logger.info(f"Frontend directory exists: {frontend_dir.exists()}")

# Serve static files for frontend
app.mount("/static", StaticFiles(directory=str(frontend_dir)), name="static")

# Serve the main HTML file at root
@app.get("/")
async def serve_frontend():
    """Serve the main frontend HTML file."""
    html_path = frontend_dir / "index.html"
    logger.info(f"Serving frontend from: {html_path}")
    logger.info(f"HTML file exists: {html_path.exists()}")
    if html_path.exists():
        return FileResponse(str(html_path))
    else:
        # Fallback if file not found
        return {
            "error": "Frontend not found",
            "message": "Please access the API at /api/info",
            "health": "/api/health",
            "frontend_path": str(html_path),
            "working_dir": str(current_dir)
        }

if __name__ == "__main__":
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    print(f"Starting server on port {port}")
    uvicorn.run(app, host="0.0.0.0", port=port)
